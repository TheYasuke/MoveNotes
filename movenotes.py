
from pptx import Presentation
import sys

#max slides
def numSlid():
	if len(oldPPT.slides) > len(newPPT.slides):
		return(len(oldPPT.slides))
	else:
		return(len(newPPT.slides))


#Loop through slides 
def slideLoop(aMatch, aItem, totalSlides, newSlideB):
	#for every slide in the powerpoint, if slide matches it will transfer notes.
	for bItem in range(totalSlides):
		try:
			if not newSlideB.shapes[aMatch].has_text_frame or not oldPPT.slides[bItem].shapes[aMatch].has_text_frame:
				continue
			if newSlideB.shapes[aMatch].text_frame.text == oldPPT.slides[bItem].shapes[aMatch].text_frame.text:
				print(f" Match slide {aItem+1} and transfer")
				newSlideB.notes_slide.notes_text_frame.text = oldPPT.slides[bItem].notes_slide.notes_text_frame.text

				break
			#runs if more old slides than new or same amount of slides and there are no matches
			elif bItem == totalSlides-1:
				print(f"No match for slide {aItem+1}")

		#runs if more new slides than old and there are no matches
		except IndexError:
			print(f"No match for slide {aItem+1}")
			break

def main():

	print ("New:", len(newPPT.slides), " Old:", len(oldPPT.slides))
	for item in range(numSlid()):
		match=0
		try: #try needed to see if item out of range for slides
			newSlide = newPPT.slides[item]
			oldSlide = oldPPT.slides[item]

			#detect if slide has title text frame if not skip and if old slide doesnt have run slide loop
			#match increases by 1, which is the determines the text frame
			if not newSlide.shapes[match].has_text_frame:
				match=1
							
			elif not oldSlide.shapes[match].has_text_frame:
				match=1
				slideLoop(match, item,numSlid(),newSlide)
				

			#if the slide numbers match transfer data, run slide loop
			if newSlide.shapes[match].text_frame.text == oldSlide.shapes[match].text_frame.text:
				print(f"Match slide {item+1} and transfer")
				newSlide.notes_slide.notes_text_frame.text = oldSlide.notes_slide.notes_text_frame.text
			elif newSlide.shapes[match].text_frame.text != oldSlide.shapes[match].text_frame.text:
				#does current new slide match next old slide
				slideLoop(match, item,numSlid(),newSlide)

		#if there are more new slides than old slides, run slide loop
		except IndexError:
			slideLoop(match, item,numSlid(),newSlide)
	try:
		newPPT.save(sys.argv[1])
	except:
		print("\nPlease close the ppt and try again.")
	else:
		print("\nTransfer complete")

try:
	
	if len(sys.argv)<3:
		print("\nComparing two similar powerpoint files, transfer notes from old powerpoint to new powerpoint.\n")
		print("If slide says 'No match' the slide may have to be updated manually.\n\nUpdatedPPT file will have notes from the oldPPT file.\n\nNote:Ensure updatedPPT is NOT open.\n\nUsage: python movenotes.py updatedPPT.pptx oldPPT.pptx")
	else:
		newPPT = Presentation(sys.argv[1])
		oldPPT = Presentation(sys.argv[2])
		main()

except Exception as e:
	print ("\nPowerpoint not found.")
	print(e)
	print ("\nUsage: python movenotes.py updatedPPT.pptx oldPPT.pptx")
	



